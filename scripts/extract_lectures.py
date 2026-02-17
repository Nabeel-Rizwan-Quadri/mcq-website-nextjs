#!/usr/bin/env python3
"""
Extract text from lecture files and store normalized output for MCQ generation.

Supported input formats:
- .pptx (requires python-pptx)
- .pdf (requires pypdf)
- .txt / .md
"""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path


SUPPORTED_SUFFIXES = {".pptx", ".pdf", ".txt", ".md"}


def slugify(value: str) -> str:
    normalized = re.sub(r"[^a-zA-Z0-9]+", "-", value.lower()).strip("-")
    return normalized or "lecture"


def normalize_whitespace(text: str) -> str:
    compact = re.sub(r"\s+", " ", text)
    return compact.strip()


def extract_from_pptx(file_path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("Missing dependency for .pptx extraction: python-pptx") from exc

    presentation = Presentation(str(file_path))
    slide_chunks: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            raw_text = getattr(shape, "text", "")
            cleaned = normalize_whitespace(raw_text)
            if cleaned:
                lines.append(cleaned)
        if lines:
            slide_chunks.append(f"Slide {slide_index}: " + " | ".join(lines))

    return "\n".join(slide_chunks)


def extract_from_pdf(file_path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("Missing dependency for .pdf extraction: pypdf") from exc

    reader = PdfReader(str(file_path))
    page_chunks: list[str] = []

    for page_index, page in enumerate(reader.pages, start=1):
        raw_text = page.extract_text() or ""
        cleaned = normalize_whitespace(raw_text)
        if cleaned:
            page_chunks.append(f"Page {page_index}: {cleaned}")

    return "\n".join(page_chunks)


def extract_from_text(file_path: Path) -> str:
    return normalize_whitespace(file_path.read_text(encoding="utf-8", errors="ignore"))


def extract_text(file_path: Path) -> str:
    if file_path.suffix == ".pptx":
        return extract_from_pptx(file_path)
    if file_path.suffix == ".pdf":
        return extract_from_pdf(file_path)
    if file_path.suffix in {".txt", ".md"}:
        return extract_from_text(file_path)
    return ""


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Extract text from lecture files.")
    parser.add_argument("--input", default="lectures", help="Input lecture directory")
    parser.add_argument("--output", default="data/lecture_texts.json", help="Output JSON path")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    input_dir = Path(args.input)
    output_path = Path(args.output)

    if not input_dir.exists():
        print(f"Input directory not found: {input_dir}")
        return 1

    lecture_files = sorted(path for path in input_dir.iterdir() if path.is_file() and path.suffix.lower() in SUPPORTED_SUFFIXES)

    if not lecture_files:
        print(f"No supported files found in {input_dir}. Supported: {', '.join(sorted(SUPPORTED_SUFFIXES))}")
        return 1

    records: list[dict[str, str]] = []
    skipped: list[str] = []

    for lecture_file in lecture_files:
        try:
            text = extract_text(lecture_file)
        except RuntimeError as exc:
            print(f"Skipped {lecture_file.name}: {exc}")
            skipped.append(lecture_file.name)
            continue

        if not text:
            skipped.append(lecture_file.name)
            continue

        records.append(
            {
                "id": slugify(lecture_file.stem),
                "title": lecture_file.stem.replace("_", " ").replace("-", " ").title(),
                "sourceFile": lecture_file.name,
                "text": text,
            }
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(records, indent=2), encoding="utf-8")

    print(f"Extracted {len(records)} lecture(s) to {output_path}")
    if skipped:
        print("Skipped files:", ", ".join(skipped))

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
